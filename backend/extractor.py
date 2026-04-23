import sys
import os
import re

def clean_text(text: str) -> str:
    """Removes extra whitespace and normalizes text for LLM pipeline."""
    if not text:
        return ""
    # Replace multiple spaces with a single space
    text = re.sub(r'[ \t]+', ' ', text)
    # Replace multiple newlines with a double newline to preserve paragraph structure
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()

def extract_pdf(file_path: str) -> str:
    """Uses PyMuPDF (fitz) with fallback to pdfplumber and OCR."""
    text = ""
    # Try PyMuPDF
    try:
        import fitz
        doc = fitz.open(file_path)
        for page in doc:
            page_text = page.get_text()
            if page_text:
                text += page_text + "\n"
        doc.close()
    except Exception as e:
        print(f"PyMuPDF failed: {e}", file=sys.stderr)
    
    if text.strip():
        return text

    # Fallback to pdfplumber
    try:
        import pdfplumber
        print("Falling back to pdfplumber...", file=sys.stderr)
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"pdfplumber failed: {e}", file=sys.stderr)

    if text.strip():
        return text

    # Fallback to OCR if scanned
    try:
        import pytesseract
        from pdf2image import convert_from_path
        print("Falling back to Tesseract OCR...", file=sys.stderr)
        images = convert_from_path(file_path)
        for img in images:
            page_text = pytesseract.image_to_string(img)
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        print(f"OCR fallback failed: {e}", file=sys.stderr)

    return text

def extract_docx(file_path: str) -> str:
    """Uses python-docx to extract text."""
    try:
        import docx
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text])
    except Exception as e:
        print(f"python-docx failed: {e}", file=sys.stderr)
        return ""

def extract_shape_text(shape):
    text = ""
    try:
        if hasattr(shape, "text") and shape.text:
            text += shape.text + "\n"
        if hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text") and cell.text:
                        text += cell.text + " "
                text += "\n"
        if hasattr(shape, "shapes"): # It's a group shape
            for s in shape.shapes:
                text += extract_shape_text(s)
    except:
        pass
    return text

def extract_pptx(file_path: str) -> str:
    """Uses python-pptx to extract text, handling tables and group shapes."""
    text = ""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                text += extract_shape_text(shape)
    except Exception as e:
        print(f"python-pptx failed: {e}", file=sys.stderr)
    return text

def extract_text(file_path: str) -> str:
    """
    Unified function that routes to the correct parser based on extension.
    """
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}", file=sys.stderr)
        return ""

    _, ext = os.path.splitext(file_path.lower())
    
    if ext == '.pdf':
        text = extract_pdf(file_path)
    elif ext == '.docx':
        text = extract_docx(file_path)
    elif ext == '.pptx':
        text = extract_pptx(file_path)
    elif ext in ['.txt', '.md', '.csv']:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
    else:
        # Fallback for unknown formats
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()

    return clean_text(text)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extractor.py <file_path>", file=sys.stderr)
        sys.exit(1)
    
    file_path = sys.argv[1]
    extracted = extract_text(file_path)
    # Output solely the extracted text to stdout so Node.js can read it
    print(extracted)
